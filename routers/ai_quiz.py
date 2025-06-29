from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Depends
from fastapi.responses import JSONResponse
from sqlalchemy.orm import Session
from datetime import datetime
import tempfile, os, json
import models  # This now refers to your unified models
from database import get_db
from dotenv import load_dotenv
import google.generativeai as genai

router = APIRouter(prefix="/ai-quizzes", tags=["AI Quiz Generation"])
load_dotenv()
# Ensure you have the API key in your .env file
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    raise ValueError("GEMINI_API_KEY not found in environment variables.")
genai.configure(api_key=GEMINI_API_KEY)

def extract_text_from_file(path):
    ext = os.path.splitext(path)[-1].lower()
    if ext == ".pdf":
        import fitz
        doc = fitz.open(path)
        return "\n".join([page.get_text() for page in doc])
    elif ext in [".docx", ".doc"]:
        import docx
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    else:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()


# --- Endpoint A: Generate questions from a file. (No DB interaction) ---
@router.post("/generate")
async def generate_ai_quiz_from_file(
    file: UploadFile = File(...),
    mcq_count: int = Form(...),
    msq_count: int = Form(...),
    marks_per_question: int = Form(...),
):
    suffix = os.path.splitext(file.filename)[-1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    try:
        content = extract_text_from_file(tmp_path)
        if not content.strip():
            raise HTTPException(status_code=400, detail="The uploaded file has no text content.")

        # Calculate total questions outside the prompt for clarity
        total_questions = mcq_count + msq_count

        prompt = f"""
        You are a master instructional designer and expert exam creator. Your task is to create a high-quality aptitude and application-based quiz from the provided text. Adherence to the exact number of questions requested is the highest priority.

        **Generation Mandate (Strict):**
        You MUST generate:
        1.  **Exactly {mcq_count}** Multiple Choice Questions (MCQ).
        2.  **Exactly {msq_count}** Multiple Select Questions (MSQ).
        3.  **One** suitable, concise title for the quiz.
        Failure to generate the precise number of questions specified is a failure of the task.

        **Question Quality Mandate:**
        - All questions must be **application-based or scenario-based**. They must test the user's ability to apply knowledge, not just recall it.
        - **Do NOT** create simple definition or fact-retrieval questions.
        - Incorrect options (distractors) must be plausible and relevant to the text's subject matter.

        **Final Verification Step:**
        - Before outputting the JSON, you must internally verify your work.
        - Confirm that the `questions` array in your JSON contains **exactly {total_questions} total objects**.
        - Confirm that there are exactly **{mcq_count} objects with `"type": "MCQ"`** and **{msq_count} objects with `"type": "MSQ"`**.

        **Output Format Constraint (Non-negotiable):**
        - Your ENTIRE response MUST be a single, raw, valid JSON object.
        - Do NOT include any introductory text, explanations, summaries, or markdown formatting like ```json.
        - Your response must begin with `{{` and end with `}}`.

        The JSON structure must be EXACTLY as follows:
        {{
            "title": "A Concise and Relevant Quiz Title",
            "questions": [
                {{
                    "question": "An application-based question text...",
                    "type": "MCQ",
                    "options": ["Plausible Option A", "Plausible Option B", "Plausible Option C", "Plausible Option D"],
                    "correct_answers": ["The exact text of the single correct option"]
                }},
                {{
                    "question": "A scenario-based question text...",
                    "type": "MSQ",
                    "options": ["Plausible Option W", "Plausible Option X", "Plausible Option Y", "Plausible Option Z"],
                    "correct_answers": ["The exact text of one correct option", "The exact text of another correct option"]
                }}
            ]
        }}

        **Source Text for Quiz Generation:**
        ---
        {content}
        """

        model = genai.GenerativeModel(model_name="gemini-1.5-flash")
        response = model.generate_content([prompt])
        
        text = response.text.strip()
        if text.startswith("```json"):
            text = text[len("```json"):].strip()
        if text.endswith("```"):
            text = text[:-3].strip()

        response_data = json.loads(text)
        
        for q in response_data.get("questions", []):
            q["marks"] = marks_per_question

        return JSONResponse(content=response_data)

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An error occurred: {str(e)}")
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


# --- Endpoint B: Save the AI-generated quiz to the UNIFIED database ---
@router.post("/submit")
async def submit_ai_generated_quiz(
    payload: dict,
    db: Session = Depends(get_db)
):
    try:
        title = payload.get("title")
        chapter_id = payload.get("chapter_id")
        questions = payload.get("questions", [])

        if not all([title, chapter_id, questions]):
            raise HTTPException(status_code=400, detail="Missing title, chapter_id, or questions.")

        new_quiz = models.Quiz(
            title=title,
            chapter_id=chapter_id,
            is_ai_generated=True,
            created_at=datetime.utcnow()
        )
        db.add(new_quiz)
        db.flush()  # <-- FIX #1: Flush here to get the new_quiz.id

        for q in questions:
            question_obj = models.Question(
                question_text=q["question"],
                quiz_id=new_quiz.id, # Now new_quiz.id is available
                marks=q.get("marks", 1),
                question_type=q["type"].lower()
            )
            db.add(question_obj)
            db.flush() # <-- FIX #2: Flush here to get the question_obj.id

            options = q.get("options", [])
            correct_answers = q.get("correct_answers", [])
            
            for option_text in options:
                is_correct = option_text in correct_answers
                option_obj = models.Option(
                    question_id=question_obj.id, # Now question_obj.id is available
                    option_text=option_text,
                    is_correct=is_correct
                )
                db.add(option_obj)

        db.commit() # Commit everything at the very end
        db.refresh(new_quiz)
        return {"msg": "AI Quiz successfully saved!", "quiz_id": new_quiz.id}
        
    except Exception as e:
        db.rollback()
        # It's helpful to log the actual error for debugging
        print(f"ERROR saving AI quiz: {e}") 
        raise HTTPException(status_code=500, detail=f"Failed to save quiz: {str(e)}")